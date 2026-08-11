"""
老六 Chat — 文件内容读取器
统一接口：read_file_content(path) → {"ok": bool, "text": str, "metadata": dict}

支持：PDF(.pdf) Word(.docx) Excel(.xlsx) PPT(.pptx) 纯文本(.txt .py .md .json .csv .log 等)
"""
import os
import io

# ---- PDF ----
try:
    from PyPDF2 import PdfReader
    _PDF_OK = True
except ImportError:
    _PDF_OK = False

# ---- Word ----
try:
    from docx import Document
    _DOCX_OK = True
except ImportError:
    _DOCX_OK = False

# ---- Excel ----
try:
    from openpyxl import load_workbook
    _XLSX_OK = True
except ImportError:
    _XLSX_OK = False

# ---- PPT ----
try:
    from pptx import Presentation
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False


def _file_size_str(path: str) -> str:
    """返回友好的文件大小字符串"""
    try:
        size = os.path.getsize(path)
        if size < 1024:
            return f"{size}B"
        elif size < 1024 * 1024:
            return f"{size / 1024:.1f}KB"
        else:
            return f"{size / (1024 * 1024):.1f}MB"
    except:
        return "?"


def _file_ext(path: str) -> str:
    return os.path.splitext(path)[1].lower()


def read_file_content(path: str) -> dict:
    """
    读取文件内容，返回统一格式。
    - ok: 是否成功
    - text: 提取的文本内容
    - metadata: {name, size, ext, pages/sheets/slides, truncated}
    - error: 错误信息（仅失败时）
    """
    if not os.path.exists(path):
        return {"ok": False, "error": "文件不存在", "metadata": {"name": os.path.basename(path)}}

    name = os.path.basename(path)
    size_str = _file_size_str(path)
    ext = _file_ext(path)
    meta = {"name": name, "size": size_str, "ext": ext}

    try:
        # ---- PDF ----
        if ext == ".pdf":
            if not _PDF_OK:
                return {"ok": False, "error": "PyPDF2 未安装（pip install PyPDF2）", "metadata": meta}
            reader = PdfReader(path)
            pages_text = []
            total = len(reader.pages)
            for i, page in enumerate(reader.pages):
                t = page.extract_text()
                if t:
                    pages_text.append(t)
                if len("\n".join(pages_text)) > 100000:
                    pages_text.append(f"\n…（PDF过大，已截取前{i+1}/{total}页）")
                    meta["truncated"] = True
                    break
            meta["pages"] = total
            return {"ok": True, "text": "\n\n".join(pages_text), "metadata": meta}

        # ---- Word (.docx) ----
        elif ext == ".docx":
            if not _DOCX_OK:
                return {"ok": False, "error": "python-docx 未安装（pip install python-docx）", "metadata": meta}
            doc = Document(path)
            paras = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n".join(paras)
            if len(text) > 100000:
                text = text[:100000] + "\n…（文档过大已截断）"
                meta["truncated"] = True
            meta["paragraphs"] = len(doc.paragraphs)
            return {"ok": True, "text": text, "metadata": meta}

        # ---- Excel (.xlsx) ----
        elif ext in (".xlsx", ".xlsm"):
            if not _XLSX_OK:
                return {"ok": False, "error": "openpyxl 未安装（pip install openpyxl）", "metadata": meta}
            wb = load_workbook(path, read_only=True, data_only=True)
            all_sheets = []
            for sname in wb.sheetnames:
                ws = wb[sname]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(v) if v is not None else "" for v in row]
                    if any(row_vals):
                        rows.append(" | ".join(row_vals))
                if rows:
                    all_sheets.append(f"## Sheet: {sname}\n" + "\n".join(rows[:200]))
                if sum(len(s) for s in all_sheets) > 100000:
                    all_sheets.append("\n…（表格过大已截断）")
                    meta["truncated"] = True
                    break
            wb.close()
            meta["sheets"] = len(wb.sheetnames)
            return {"ok": True, "text": "\n\n".join(all_sheets), "metadata": meta}

        # ---- PPT (.pptx) ----
        elif ext == ".pptx":
            if not _PPTX_OK:
                return {"ok": False, "error": "python-pptx 未安装（pip install python-pptx）", "metadata": meta}
            prs = Presentation(path)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                lines = [f"--- 第{i+1}页 ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)
                if len(lines) > 1:
                    slides_text.append("\n".join(lines))
            text = "\n\n".join(slides_text)
            if len(text) > 100000:
                text = text[:100000] + "\n…（PPT过大已截断）"
                meta["truncated"] = True
            meta["slides"] = len(prs.slides)
            return {"ok": True, "text": text, "metadata": meta}

        # ---- 纯文本文件 ----
        else:
            # 检测是否为二进制
            with open(path, "rb") as f:
                head = f.read(1024)
            if b"\x00" in head:
                return {"ok": False, "error": "无法读取的二进制文件", "metadata": meta}

            with open(path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            if len(text) > 50000:
                text = text[:50000] + f"\n…（文件过大已截断，原文{len(text)}字符）"
                meta["truncated"] = True
            meta["chars"] = len(text)
            return {"ok": True, "text": text, "metadata": meta}

    except Exception as e:
        return {"ok": False, "error": str(e), "metadata": meta}


def get_file_type_icon(ext: str) -> str:
    """根据扩展名返回对应的 emoji 图标"""
    icons = {
        ".pdf": "📕",
        ".docx": "📘", ".doc": "📘",
        ".xlsx": "📗", ".xls": "📗", ".csv": "📗",
        ".pptx": "📙", ".ppt": "📙",
        ".py": "🐍",
        ".js": "📜", ".ts": "📜",
        ".json": "📋",
        ".txt": "📄", ".md": "📝",
        ".png": "🖼", ".jpg": "🖼", ".jpeg": "🖼", ".gif": "🖼", ".bmp": "🖼", ".webp": "🖼",
        ".zip": "📦", ".rar": "📦", ".7z": "📦",
    }
    return icons.get(ext, "📄")


def is_supported(ext: str) -> bool:
    """检查文件类型是否受支持"""
    supported = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx",
                 ".txt", ".py", ".md", ".json", ".csv", ".log",
                 ".js", ".ts", ".html", ".css", ".xml", ".yaml", ".yml",
                 ".ini", ".cfg", ".toml", ".sh", ".bat", ".ps1"}
    return ext in supported
